"""
Document exporters — slide deck, PDF report and templated manuscript.

All generated locally from content the pipeline already produced (sections,
papers, synthesis, comparison table, year distribution). No LLM calls here,
so exporting costs nothing.

Formats:
  build_pptx()  -> slide deck (.pptx) — includes a real data table and a
                   native chart slide, not just bullet points.
  build_pdf()   -> formatted report (.pdf) — mirrors the in-app Review tab
                   exactly (written sections + a plain reference list).
  build_docx()  -> manuscript in an IEEE / arXiv style template (.docx).

Citations: references are plain, standard bibliographic text (no
hyperlinks) — the inline "[n]" markers within the written prose are the
clickable part, each linking to that paper's URL.

Text sanitization: paper metadata pulled from external APIs sometimes carries
Unicode punctuation (odd dash/hyphen variants, smart quotes) that the core
PDF/DOCX fonts can't render, showing up as black-box "missing glyph" marks.
`_sanitize()` normalizes that away before anything is written out.
"""
import io
import re
import unicodedata
from datetime import date

SECTION_ORDER = [
    ("abstract", "Abstract"),
    ("intro", "Introduction"),
    ("synthesis", "Literature Review"),
    ("gaps", "Gaps and Limitations"),
    ("future", "Future Directions and Conclusion"),
]

ACCENT = (0x5B, 0x4F, 0xF0)
LINK_COLOR = (0x1A, 0x5F, 0xD9)


# ── text sanitization ────────────────────────────────────────────────────────

# Various dash/hyphen/minus lookalikes that fall outside the core PDF/DOCX
# fonts' supported glyph set and render as a black box — fold them all to a
# plain hyphen. (U+2010..U+2015 dash forms, U+2212 minus sign.)
_DASH_CHARS = "‐‑‒–—―−"
_QUOTE_MAP = {
    "‘": "'", "’": "'", "‚": "'", "‛": "'", "ʼ": "'", "´": "'", "′": "'", "`": "'",
    "“": '"', "”": '"', "„": '"', "‟": '"',
}
_TAG_RE = re.compile(r"</?[a-zA-Z][a-zA-Z0-9]*(?:\s[^<>]*)?>")


def _sanitize(text) -> str:
    """Normalize Unicode punctuation that commonly shows up as a missing-
    glyph box in exported PDFs/DOCX (odd dashes, smart quotes, stray
    replacement characters from an upstream encoding hiccup), and strip stray
    HTML/XML markup (e.g. "<scp>A</scp>lzheimer's disease") that some source
    APIs (Semantic Scholar / OpenAlex) leak straight through in titles."""
    if not text:
        return ""
    t = unicodedata.normalize("NFKC", str(text))
    t = _TAG_RE.sub("", t)
    for d in _DASH_CHARS:
        t = t.replace(d, "-")
    for k, v in _QUOTE_MAP.items():
        t = t.replace(k, v)
    t = t.replace("�", "")  # the literal "unknown character" replacement glyph
    return "".join(ch for ch in t if ch == "\n" or ch.isprintable())


# ── inline [n] citation linking ─────────────────────────────────────────────
# References themselves stay plain, standard bibliographic text (no
# hyperlinks) — that's the convention. Instead, the inline "[n]" citation
# markers that already appear throughout the written prose are what get
# turned into clickable links, each pointing at that specific paper's URL.

_CITE_RE = re.compile(r"(\[\d+\])")


def _split_citations(text: str) -> list[tuple[str, int | None]]:
    """Split text into (segment, citation_number) pairs. citation_number is
    None for a plain-text segment, or the int inside a "[n]" marker."""
    out = []
    for part in _CITE_RE.split(text or ""):
        if not part:
            continue
        m = re.fullmatch(r"\[(\d+)\]", part)
        out.append((part, int(m.group(1)) if m else None))
    return out


def citation_urls(papers: list[dict]) -> dict[int, str]:
    """Map citation number (1-based, matching reference_list()'s numbering)
    -> that paper's URL, for linking inline [n] markers."""
    out = {}
    for i, p in enumerate(papers or [], 1):
        url = (p.get("url") or "").strip()
        if url:
            out[i] = url
    return out


def _paras(text: str) -> list[str]:
    return [_sanitize(p.strip()) for p in re.split(r"\n\s*\n", text or "") if p.strip()]


def review_title(sections: dict, topic: str) -> str:
    return _sanitize((sections or {}).get("title") or (topic or "Literature Review"))


def reference_list(papers: list[dict]) -> list[dict]:
    """IEEE-style numbered references in citation order, as
    {"text": "...", "url": "..."} so exporters can make them clickable."""
    out = []
    for i, p in enumerate(papers or [], 1):
        authors = _sanitize(p.get("authors")) or "—"
        title = _sanitize(p.get("title"))
        venue = _sanitize(p.get("venue")) or "preprint"
        year = p.get("year") or "n.d."
        out.append({
            "text": f'[{i}] {authors}, "{title}," {venue}, {year}.',
            "url": (p.get("url") or "").strip(),
        })
    return out


# ── PPTX ───────────────────────────────────────────────────────────────────

def build_pptx(topic: str, sections: dict, papers: list[dict], synthesis: dict,
                comparison: list[dict] | None = None, year_dist: list[dict] | None = None,
                images: dict[str, bytes] | None = None,
                slide_bullets: dict[str, list[str]] | None = None) -> bytes:
    """`images` (optional, opt-in, costs money — see core/image_gen.py) maps a
    SECTION_ORDER key ("intro", "synthesis", ...) to a generated illustration;
    that section's FIRST slide gets a narrower bullet column plus the image,
    everything else is unchanged when omitted (the normal, free export).

    `slide_bullets` (optional — see agents/slide_writer.py) maps a section key
    to presentation-native bullets for that section, replacing the default
    mechanical sentence-split of the report prose. A section missing from
    this dict (call failed, or the caller didn't generate it) just falls back
    to the old behavior, so this is always safe to omit or partially fill."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)   # 16:9
    accent = RGBColor(*ACCENT)
    link_rgb = RGBColor(*LINK_COLOR)
    accent_light = RGBColor(*(min(255, c + 130) for c in ACCENT))
    accent_mid = RGBColor(*(min(255, c + 65) for c in ACCENT))
    title = review_title(sections, topic)
    cite_urls = citation_urls(papers)

    def _no_line(shape):
        shape.line.fill.background()

    def _no_dashed(shape):
        """Dashed outline, no fill — used for the 'gaps' motif to visually
        read as open/incomplete rather than solid."""
        shape.fill.background()
        shape.line.color.rgb = accent
        shape.line.width = Pt(2.5)
        ln = shape.line._get_or_add_ln()
        dash = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
        ln.append(dash)

    def draw_motif(s, key):
        """A small set of free, programmatic vector diagrams — one per
        section type — drawn with plain pptx autoshapes. No AI model, no
        cost, always available (unlike the opt-in AI illustration below,
        which overrides this when the caller supplies a real image)."""
        x, y, size = Inches(9.3), Inches(2.6), Inches(3.2)
        try:
            if key == "abstract":
                # a ring — "the whole picture, at a glance"
                shp = s.shapes.add_shape(MSO_SHAPE.DONUT, x, y, size, size)
                shp.fill.solid(); shp.fill.fore_color.rgb = accent_light
                _no_line(shp)
            elif key == "intro":
                # narrowing bars — scoping down into the topic
                widths = [size, size * 0.72, size * 0.46]
                colors = [accent_light, accent_mid, accent]
                for i, (w, col) in enumerate(zip(widths, colors)):
                    bar = s.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        x + (size - w) / 2, y + Inches(0.05) + i * Inches(1.05), w, Inches(0.75))
                    bar.fill.solid(); bar.fill.fore_color.rgb = col
                    _no_line(bar)
            elif key == "synthesis":
                # three overlapping circles — themes converging
                r = size * 0.62
                offs = [(0, 0, accent_light), (size - r, 0, accent_mid), (size / 2 - r / 2, size - r, accent)]
                for dx, dy, col in offs:
                    c = s.shapes.add_shape(MSO_SHAPE.OVAL, x + dx, y + dy, r, r)
                    c.fill.solid(); c.fill.fore_color.rgb = col
                    _no_line(c)
            elif key == "gaps":
                # a dashed, unfilled ring — something open/incomplete
                shp = s.shapes.add_shape(MSO_SHAPE.OVAL, x + size * 0.1, y + size * 0.1,
                                          size * 0.8, size * 0.8)
                _no_dashed(shp)
            elif key == "future":
                # chevrons pointing forward — next steps
                for i in range(3):
                    w = size * 0.42
                    chev = s.shapes.add_shape(
                        MSO_SHAPE.CHEVRON, x + i * (w * 0.72), y + size * 0.35, w, size * 0.3)
                    chev.fill.solid()
                    chev.fill.fore_color.rgb = [accent_light, accent_mid, accent][i]
                    _no_line(chev)
        except Exception:
            pass  # a motif is a nice-to-have, never worth failing the export over

    def bullets_slide(heading, bullets, note=None, linkify=False, image_bytes=None, motif_key=None):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = heading
        s.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(30)
        s.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = accent
        body_ph = s.placeholders[1]
        if image_bytes:
            # Narrow the bullet column to make room for the illustration on
            # the right instead of the usual full-width text block.
            body_ph.left, body_ph.top = Inches(0.5), Inches(1.6)
            body_ph.width, body_ph.height = Inches(6.9), Inches(5.4)
            try:
                s.shapes.add_picture(io.BytesIO(image_bytes), Inches(7.8), Inches(1.6),
                                      width=Inches(5.0), height=Inches(5.0))
            except Exception:
                pass  # a corrupt/unsupported image should never break the export
        elif motif_key:
            # Free vector fallback/default — narrower column, smaller motif
            # (it's an accent, not a full illustration).
            body_ph.left, body_ph.top = Inches(0.5), Inches(1.6)
            body_ph.width, body_ph.height = Inches(8.3), Inches(5.4)
            draw_motif(s, motif_key)
        body = body_ph.text_frame
        body.word_wrap = True
        first = True
        for b in bullets:
            p = body.paragraphs[0] if first else body.add_paragraph()
            segs = _split_citations(b) if linkify else [(b, None)]
            for seg, num in segs:
                run = p.add_run()
                run.text = seg
                run.font.size = Pt(16)
                if num and cite_urls.get(num):
                    run.hyperlink.address = cite_urls[num]
                    run.font.color.rgb = link_rgb
                    run.font.underline = True
            first = False
        if note:
            s.notes_slide.notes_text_frame.text = note
        return s

    def blank_titled_slide(heading):
        s = prs.slides.add_slide(prs.slide_layouts[5])  # title-only layout
        s.shapes.title.text = heading
        s.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
        s.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = accent
        return s

    # Title slide. Generated review titles run long (full academic-style
    # titles, often 15-20 words) — the default Title Slide layout's
    # placeholder is sized for a short deck name and a long title overflows
    # it, bleeding down into the subtitle placeholder underneath (visible as
    # garbled overlapping text). Give the title its own generously-sized,
    # explicitly-positioned box, scale the font down for longer titles, and
    # push the subtitle well clear of it instead of trusting the layout's
    # fixed default geometry.
    from pptx.enum.text import PP_ALIGN

    s0 = prs.slides.add_slide(prs.slide_layouts[0])
    title_ph = s0.shapes.title
    title_ph.left, title_ph.top = Inches(1.0), Inches(1.5)
    title_ph.width, title_ph.height = Inches(11.3), Inches(3.0)
    title_ph.text = title
    tf = title_ph.text_frame
    tf.word_wrap = True
    title_size = 34 if len(title) <= 70 else (28 if len(title) <= 110 else 24)
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(title_size)
            run.font.bold = True
            run.font.color.rgb = accent

    subtitle_ph = s0.placeholders[1]
    subtitle_ph.left, subtitle_ph.top = Inches(1.0), Inches(4.9)
    subtitle_ph.width, subtitle_ph.height = Inches(11.3), Inches(1.6)
    subtitle_ph.text = (
        f"A literature review of {len(papers)} sources\n"
        f"Generated by Sift · {date.today().isoformat()}"
    )
    for p in subtitle_ph.text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER

    # Content slides from each written section
    for key, label in SECTION_ORDER:
        text = (sections or {}).get(key)
        if not text:
            continue
        # Prefer slide-native bullets (agents/slide_writer.py, Gemini) when the
        # caller supplied them — actual presentation-style points instead of
        # academic sentences mechanically chopped out of the report prose.
        # Falls back to the old sentence-split for any section that wasn't
        # rewritten (call failed, rate-limited, or the caller opted out).
        written = [_sanitize(b) for b in (slide_bullets or {}).get(key, []) if (b or "").strip()]
        if written:
            bullets = written
        else:
            chunks = _paras(text)
            # split long sections across slides (max ~5 bullets each)
            bullets = []
            for c in chunks:
                for sent in re.split(r"(?<=[.!?])\s+", c):
                    sent = sent.strip()
                    if len(sent) > 12:
                        bullets.append(sent if len(sent) <= 220 else sent[:217] + "…")
        section_image = (images or {}).get(key)
        for n in range(0, len(bullets), 5):
            part = bullets[n:n + 5]
            head = label if n == 0 else f"{label} (cont.)"
            bullets_slide(
                head, part, linkify=True,
                image_bytes=section_image if n == 0 else None,
                # No motif by default — the abstract shapes (ring, bars,
                # circles...) had no label explaining what they represented,
                # so they just read as unexplained decoration rather than
                # content. Bullets get the full slide width instead. An
                # explicit AI illustration (opt-in, costs money — see
                # ExportBar.jsx) still takes priority when present.
                motif_key=None,
            )

    # Themes & gaps from the synthesis
    themes = [_sanitize(t) for t in ((synthesis or {}).get("themes") or [])]
    gaps = [_sanitize(g) for g in ((synthesis or {}).get("gaps") or [])]
    if themes:
        bullets_slide("Key Themes", themes[:6])
    if gaps:
        bullets_slide("Research Gaps", gaps[:6])

    # ── Data: publication-year chart (a real native chart, not a bullet list) ──
    if year_dist and len(year_dist) > 1:
        s = blank_titled_slide("Papers by Publication Year")
        chart_data = CategoryChartData()
        chart_data.categories = [str(d.get("year", "?")) for d in year_dist]
        chart_data.add_series("Papers", [d.get("count", 0) for d in year_dist])
        x, y, cx, cy = Inches(1.0), Inches(1.6), Inches(11.3), Inches(5.4)
        gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
        chart = gframe.chart
        chart.has_legend = False
        try:
            plot = chart.plots[0]
            plot.series[0].format.fill.solid()
            plot.series[0].format.fill.fore_color.rgb = accent
        except Exception:
            pass

    # ── Data: comparison table (real pptx table, not bullets) ──
    if comparison:
        rows_per_slide = 8
        cols = ["#", "Paper", "Year", "Method", "Key finding"]
        for start in range(0, len(comparison), rows_per_slide):
            chunk = comparison[start:start + rows_per_slide]
            label = "Paper Comparison" if start == 0 else "Paper Comparison (cont.)"
            s = blank_titled_slide(label)
            nrows, ncols = len(chunk) + 1, len(cols)
            x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5)
            tbl_shape = s.shapes.add_table(nrows, ncols, x, y, cx, cy)
            tbl = tbl_shape.table
            widths = [Inches(0.5), Inches(4.6), Inches(0.9), Inches(2.8), Inches(3.5)]
            for i, w in enumerate(widths):
                tbl.columns[i].width = w
            for c, head in enumerate(cols):
                cell = tbl.cell(0, c)
                cell.text = head
                cell.text_frame.paragraphs[0].runs[0].font.bold = True
                cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
                cell.fill.solid()
                cell.fill.fore_color.rgb = accent
                cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            for r, row in enumerate(chunk, start=1):
                vals = [
                    str(start + r),  # display rank (1-based position), not the raw paper idx
                    _sanitize(row.get("title"))[:90],
                    str(row.get("year") or "—"),
                    _sanitize(row.get("method"))[:60] or "—",
                    _sanitize(row.get("finding"))[:90] or "—",
                ]
                for c, v in enumerate(vals):
                    cell = tbl.cell(r, c)
                    cell.text = v
                    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(10.5)

    # References — standard bibliographic text, not hyperlinked (the inline
    # [n] markers in the section slides above carry the links instead).
    refs = reference_list(papers)
    for n in range(0, len(refs), 8):
        chunk = refs[n:n + 8]
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = "References" if n == 0 else "References (cont.)"
        s.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(30)
        s.shapes.title.text_frame.paragraphs[0].runs[0].font.color.rgb = accent
        body = s.placeholders[1].text_frame
        body.word_wrap = True
        first = True
        for ref in chunk:
            p = body.paragraphs[0] if first else body.add_paragraph()
            p.text = ref["text"]
            p.font.size = Pt(14)
            first = False

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── PDF ────────────────────────────────────────────────────────────────────

def build_pdf(topic: str, sections: dict, papers: list[dict], synthesis: dict,
              comparison: list[dict] | None = None, year_dist: list[dict] | None = None) -> bytes:
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=LETTER, title=review_title(sections, topic),
        leftMargin=0.9 * inch, rightMargin=0.9 * inch,
        topMargin=0.85 * inch, bottomMargin=0.85 * inch,
    )
    ss = getSampleStyleSheet()
    accent = colors.Color(ACCENT[0] / 255, ACCENT[1] / 255, ACCENT[2] / 255)
    link_color = colors.Color(LINK_COLOR[0] / 255, LINK_COLOR[1] / 255, LINK_COLOR[2] / 255)
    h_title = ParagraphStyle("t", parent=ss["Title"], fontSize=19, leading=24, textColor=accent)
    h_sec = ParagraphStyle("s", parent=ss["Heading2"], fontSize=13, leading=17,
                           spaceBefore=14, spaceAfter=6, textColor=accent)
    body = ParagraphStyle("b", parent=ss["BodyText"], fontSize=10.5, leading=15.5,
                          alignment=4, spaceAfter=8)        # justified
    small = ParagraphStyle("m", parent=ss["BodyText"], fontSize=9, textColor=colors.grey)
    ref = ParagraphStyle("r", parent=body, leftIndent=18, firstLineIndent=-18, spaceAfter=4)

    def esc(t):
        return (_sanitize(t) or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    link_hex = f"{LINK_COLOR[0]:02x}{LINK_COLOR[1]:02x}{LINK_COLOR[2]:02x}"
    cite_urls = citation_urls(papers)

    def linkify(text: str) -> str:
        """Escape body text, then turn inline [n] citation markers into
        clickable links to that paper's URL (reportlab supports inline <a>
        markup directly in Paragraph text)."""
        out = []
        for seg, num in _split_citations(text):
            e = esc(seg)
            if num and cite_urls.get(num):
                safe_url = cite_urls[num].replace("&", "&amp;")
                e = f'<a href="{safe_url}" color="#{link_hex}"><u>{e}</u></a>'
            out.append(e)
        return "".join(out)

    story = [Paragraph(esc(review_title(sections, topic)), h_title)]
    story.append(Paragraph(
        f"A literature review of {len(papers)} sources · generated by Sift · {date.today().isoformat()}",
        small))
    story.append(Spacer(1, 6))
    story.append(HRFlowable(width="100%", color=accent, thickness=1))

    # This PDF mirrors the in-app Review tab exactly: the written sections,
    # then a plain reference list. No extra data table here — that lives in
    # the slide deck, where it adds real value as its own visual.
    for key, label in SECTION_ORDER:
        text = (sections or {}).get(key)
        if not text:
            continue
        story.append(Paragraph(label, h_sec))
        for p in _paras(text):
            story.append(Paragraph(linkify(p), body))

    if papers:
        story.append(Paragraph("References", h_sec))
        for r in reference_list(papers):
            story.append(Paragraph(esc(r["text"]), ref))

    doc.build(story)
    return buf.getvalue()


# ── DOCX (IEEE / arXiv templates) ──────────────────────────────────────────

def _add_hyperlink(paragraph, url: str, text: str, size_pt: int | None = None):
    """python-docx has no high-level hyperlink API — this is the standard
    OOXML recipe: register the URL as an external relationship on the part,
    then build the <w:hyperlink> run manually so it's a real clickable link
    (not just blue underlined text)."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    from docx.opc.constants import RELATIONSHIP_TYPE

    part = paragraph.part
    r_id = part.relate_to(url, RELATIONSHIP_TYPE.HYPERLINK, is_external=True)

    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), r_id)

    new_run = OxmlElement("w:r")
    rpr = OxmlElement("w:rPr")
    color = OxmlElement("w:color")
    color.set(qn("w:val"), f"{LINK_COLOR[0]:02X}{LINK_COLOR[1]:02X}{LINK_COLOR[2]:02X}")
    rpr.append(color)
    u = OxmlElement("w:u")
    u.set(qn("w:val"), "single")
    rpr.append(u)
    if size_pt:
        sz = OxmlElement("w:sz")
        sz.set(qn("w:val"), str(size_pt * 2))
        rpr.append(sz)
    new_run.append(rpr)
    t = OxmlElement("w:t")
    t.text = text
    new_run.append(t)
    hyperlink.append(new_run)
    paragraph._p.append(hyperlink)
    return hyperlink


def build_docx(topic: str, sections: dict, papers: list[dict], synthesis: dict,
               template: str = "ieee",
               comparison: list[dict] | None = None, year_dist: list[dict] | None = None) -> bytes:
    from docx import Document
    from docx.enum.section import WD_SECTION
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt, Inches, RGBColor

    tpl = (template or "ieee").lower()
    doc = Document()
    title = review_title(sections, topic)

    # Page + base font per template
    sec = doc.sections[0]
    for m in ("top_margin", "bottom_margin", "left_margin", "right_margin"):
        setattr(sec, m, Inches(0.75 if tpl == "ieee" else 1.0))
    normal = doc.styles["Normal"]
    normal.font.name = "Times New Roman"
    normal.font.size = Pt(10 if tpl == "ieee" else 11)

    # ── Title block (single column, even in IEEE) ──
    t = doc.add_paragraph()
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = t.add_run(title)
    run.font.size = Pt(20 if tpl == "ieee" else 17)
    run.font.bold = True

    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    mrun = meta.add_run(f"Generated by Sift · {len(papers)} sources · {date.today().isoformat()}")
    mrun.font.size = Pt(9)
    mrun.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # IEEE is two-column from the abstract onward
    if tpl == "ieee":
        new = doc.add_section(WD_SECTION.CONTINUOUS)
        cols = new._sectPr.xpath("./w:cols")[0]
        cols.set(_qn("w:num"), "2")
        cols.set(_qn("w:space"), "360")          # ~0.25"
        for m in ("top_margin", "bottom_margin", "left_margin", "right_margin"):
            setattr(new, m, Inches(0.75))

    def heading(text, n=1):
        h = doc.add_paragraph()
        r = h.add_run(text.upper() if tpl == "ieee" else text)
        r.font.bold = True
        r.font.size = Pt(10 if tpl == "ieee" else 13)
        if tpl == "ieee":
            h.alignment = WD_ALIGN_PARAGRAPH.CENTER
        h.paragraph_format.space_before = Pt(10)
        h.paragraph_format.space_after = Pt(4)

    cite_urls = citation_urls(papers)

    def para(text, italic=False, bold_lead=None, linkify=False):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        p.paragraph_format.first_line_indent = Inches(0.2 if tpl == "ieee" else 0)
        p.paragraph_format.space_after = Pt(6)
        if bold_lead:
            b = p.add_run(bold_lead)
            b.font.bold = True
            b.font.italic = True
        clean = _sanitize(text)
        if linkify:
            # Turn inline [n] citation markers into real hyperlinks to that
            # paper's URL; everything else is one plain italic-matching run.
            for seg, num in _split_citations(clean):
                if num and cite_urls.get(num):
                    _add_hyperlink(p, cite_urls[num], seg)
                else:
                    r = p.add_run(seg)
                    r.font.italic = italic
        else:
            r = p.add_run(clean)
            r.font.italic = italic
        return p

    # ── Abstract ──
    if sections.get("abstract"):
        for i, blk in enumerate(_paras(sections["abstract"])):
            para(blk, italic=(tpl == "ieee"), bold_lead="Abstract— " if (tpl == "ieee" and i == 0) else None)

    # Index terms (IEEE convention) from synthesis themes
    themes = [_sanitize(t) for t in ((synthesis or {}).get("themes") or [])][:5]
    if themes and tpl == "ieee":
        para(", ".join(themes) + ".", italic=True, bold_lead="Index Terms— ")

    # ── Body sections ──
    numbered = ["I.", "II.", "III.", "IV.", "V.", "VI."]
    n = 0
    for key, label in SECTION_ORDER:
        if key == "abstract" or not sections.get(key):
            continue
        prefix = f"{numbered[n]} " if (tpl == "ieee" and n < len(numbered)) else ""
        heading(prefix + label)
        n += 1
        for blk in _paras(sections[key]):
            para(blk, linkify=True)

    # ── Data: comparison table ──
    if comparison:
        heading(f"{numbered[n]} Paper Comparison" if (tpl == "ieee" and n < len(numbered)) else "Paper Comparison")
        n += 1
        tbl = doc.add_table(rows=1, cols=5)
        tbl.style = "Light Grid Accent 1" if "Light Grid Accent 1" in [s.name for s in doc.styles] else "Table Grid"
        hdr = tbl.rows[0].cells
        for c, label in enumerate(["#", "Paper", "Year", "Method", "Key finding"]):
            hdr[c].text = label
            for p in hdr[c].paragraphs:
                for r in p.runs:
                    r.font.bold = True
                    r.font.size = Pt(8)
        for i, row in enumerate(comparison, 1):
            cells = tbl.add_row().cells
            vals = [str(i), _sanitize(row.get("title"))[:120], str(row.get("year") or "—"),
                    _sanitize(row.get("method"))[:80] or "—", _sanitize(row.get("finding"))[:100] or "—"]
            for c, v in enumerate(vals):
                cells[c].text = v
                for p in cells[c].paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(8)

    # ── References — plain, standard bibliographic text (not hyperlinked;
    # the inline [n] markers in the body sections above carry the links) ──
    if papers:
        heading(f"{numbered[n]} References" if (tpl == "ieee" and n < len(numbered)) else "References")
        for r in reference_list(papers):
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            p.paragraph_format.left_indent = Inches(0.25)
            p.paragraph_format.first_line_indent = Inches(-0.25)
            p.paragraph_format.space_after = Pt(3)
            size = 8.5 if tpl == "ieee" else 10
            run = p.add_run(r["text"])
            run.font.size = Pt(size)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _qn(tag: str):
    from docx.oxml.ns import qn
    return qn(tag)


# ── Experiment plan exports (Methods tab) ───────────────────────────────────
# Deliberately simpler than build_docx/build_pdf above — no IEEE/arXiv two-
# column manuscript layout, since this is a working document a researcher
# edits, prints for a lab meeting, or attaches to a proposal draft, not a
# submission template. Both functions take the same shape of data the
# frontend already renders in MethodsPanel.jsx: the plan from
# ExperimentDesignerAgent, the source papers/extractions (for the evidence
# trail), and the critic's scores if a Refine pass has run.

def _evidence_items(h: dict) -> tuple[list[dict], list[str]]:
    """Split a hypothesis's approaches+baselines into (grounded, proposed) —
    same split MethodsPanel.jsx's EvidenceTrail does, so the export matches
    what's on screen."""
    items = [{"name": a.get("name"), "from_idx": a.get("from_idx")}
             for a in (h.get("approaches") or [])] + \
            [{"name": b.get("name"), "from_idx": b.get("from_idx")}
             for b in (h.get("baselines") or [])]
    grounded = [it for it in items if it.get("from_idx") is not None]
    proposed = [it["name"] for it in items if it.get("from_idx") is None]
    return grounded, proposed


def build_experiments_docx(topic: str, plan: dict, papers: list[dict],
                            extractions: list[dict] | None = None,
                            critique: dict | None = None) -> bytes:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt, Inches, RGBColor

    doc = Document()
    sec = doc.sections[0]
    for m in ("top_margin", "bottom_margin", "left_margin", "right_margin"):
        setattr(sec, m, Inches(1.0))
    normal = doc.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(11)

    title_p = doc.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title_p.add_run(_sanitize(topic) or "Experiment Plan")
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor(*ACCENT)

    n = len(papers or [])
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    mrun = meta.add_run(
        f"Generated by Sift's Method & Experiment Designer · grounded in "
        f"{n} source{'s' if n != 1 else ''} · {date.today().isoformat()}"
    )
    mrun.font.size = Pt(9)
    mrun.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    if plan.get("domain"):
        p = doc.add_paragraph()
        p.add_run("Domain: ").bold = True
        p.add_run(_sanitize(plan["domain"]))
    if plan.get("note"):
        p = doc.add_paragraph()
        r = p.add_run(_sanitize(plan["note"]))
        r.font.italic = True
        r.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    ext_by_idx = {e.get("idx"): e for e in (extractions or [])}
    title_by_idx = {p.get("idx"): p.get("title") for p in (papers or [])}
    crit_by_index = {c.get("index"): c for c in (critique or {}).get("critiques", [])}

    def field(label, text):
        if not text:
            return
        p = doc.add_paragraph()
        p.add_run(f"{label}: ").bold = True
        p.add_run(_sanitize(text))

    def table(headers, rows):
        t = doc.add_table(rows=1, cols=len(headers))
        t.style = "Light Grid Accent 1"
        for j, label in enumerate(headers):
            t.rows[0].cells[j].text = label
        for row_vals in rows:
            cells = t.add_row().cells
            for j, v in enumerate(row_vals):
                cells[j].text = "" if v is None else str(v)

    for i, h in enumerate(plan.get("hypotheses") or []):
        doc.add_paragraph()  # spacer
        hd = doc.add_paragraph()
        r = hd.add_run(f"H{i + 1}. {_sanitize(h.get('hypothesis'))}")
        r.font.bold = True
        r.font.size = Pt(13)
        r.font.color.rgb = RGBColor(*ACCENT)

        if h.get("rationale"):
            p = doc.add_paragraph()
            r = p.add_run(_sanitize(h["rationale"]))
            r.font.italic = True
            r.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

        c = crit_by_index.get(i)
        if c:
            s = c.get("scores") or {}
            table(["Novelty", "Grounding", "Testability", "Consistency", "Overall"],
                  [[s.get("novelty", "—"), s.get("grounding", "—"), s.get("testability", "—"),
                    s.get("consistency", "—"), c.get("overall", "—")]])
            if c.get("issues"):
                doc.add_paragraph("Critic: " + "; ".join(_sanitize(x) for x in c["issues"]))

        field("Setup", h.get("setup"))
        v = h.get("variables") or {}
        if v:
            field("Variables",
                  f"independent: {v.get('independent', '')} · dependent: {v.get('dependent', '')} "
                  f"· controlled: {v.get('controlled', '')}")

        grounded, proposed = _evidence_items(h)
        if grounded:
            doc.add_paragraph().add_run("Evidence trail").bold = True
            table(["Approach / baseline", "Source paper", "What it found"], [
                [_sanitize(it["name"]),
                 _sanitize(title_by_idx.get(it["from_idx"])) or f"[{it['from_idx']}]",
                 _sanitize((ext_by_idx.get(it["from_idx"]) or {}).get("finding")) or "—"]
                for it in grounded
            ])
        if proposed:
            doc.add_paragraph(
                "Not traced to a source (the agent's own proposal): "
                + ", ".join(_sanitize(x) for x in proposed))

        metrics = h.get("metrics") or []
        if metrics:
            doc.add_paragraph().add_run("Metrics").bold = True
            table(["Metric", "Unit / scale", "Target"], [
                [_sanitize(m.get("name")), _sanitize(m.get("unit")) or "—", _sanitize(m.get("target")) or "—"]
                for m in metrics
            ])

        if h.get("failure_modes"):
            field("What could invalidate it", "; ".join(h["failure_modes"]))
        field("Validation", h.get("validation"))
        field("Risks / ethics", h.get("risks"))

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def build_experiments_pdf(topic: str, plan: dict, papers: list[dict],
                           extractions: list[dict] | None = None,
                           critique: dict | None = None) -> bytes:
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                     HRFlowable, Table, TableStyle)

    buf = io.BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=LETTER, title=_sanitize(topic) or "Experiment Plan",
        leftMargin=0.9 * inch, rightMargin=0.9 * inch,
        topMargin=0.85 * inch, bottomMargin=0.85 * inch,
    )
    ss = getSampleStyleSheet()
    accent = colors.Color(ACCENT[0] / 255, ACCENT[1] / 255, ACCENT[2] / 255)
    h_title = ParagraphStyle("t", parent=ss["Title"], fontSize=19, leading=24, textColor=accent)
    h_hyp = ParagraphStyle("h", parent=ss["Heading2"], fontSize=13, leading=17,
                            spaceBefore=14, spaceAfter=4, textColor=accent)
    body = ParagraphStyle("b", parent=ss["BodyText"], fontSize=10.5, leading=15, spaceAfter=6)
    small = ParagraphStyle("m", parent=ss["BodyText"], fontSize=9, textColor=colors.grey)
    italic = ParagraphStyle("i", parent=body, fontName="Helvetica-Oblique", textColor=colors.grey)
    bold_lbl = ParagraphStyle("bl", parent=body, fontName="Helvetica-Bold")

    def esc(t):
        return (_sanitize(t) or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    th_style = ParagraphStyle("th", parent=ss["BodyText"], fontSize=8.5, leading=11,
                               fontName="Helvetica-Bold", textColor=colors.black)
    td_style = ParagraphStyle("td", parent=ss["BodyText"], fontSize=8.5, leading=11,
                               textColor=colors.black)

    def table(headers, rows, col_widths=None):
        # reportlab's Table does NOT word-wrap plain strings — a cell that's
        # just a str is drawn at its natural width and overflows into
        # neighboring cells/rows for anything longer than a couple of words
        # (this is what produced the garbled/overlapping export). Wrapping
        # every cell in a Paragraph makes it a real flowable that wraps to
        # the column width instead, same as every other table below.
        data = [[Paragraph(str(v), th_style) for v in headers]]
        for row_vals in rows:
            data.append([Paragraph("—" if v is None else str(v), td_style) for v in row_vals])
        # A fixed width is required for wrapping to have something to wrap
        # against — without colWidths, Table sizes columns from content and
        # a Paragraph reports its *minimum* width, which collapses columns
        # back to the same overflow problem. Split the usable page width
        # evenly when the caller doesn't specify explicit widths.
        usable_width = LETTER[0] - 1.8 * inch  # left + right margin above
        widths = col_widths or [usable_width / len(headers)] * len(headers)
        t = Table(data, colWidths=widths, hAlign="LEFT")
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.Color(0.95, 0.95, 0.97)),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.lightgrey),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ("LEFTPADDING", (0, 0), (-1, -1), 6),
            ("RIGHTPADDING", (0, 0), (-1, -1), 6),
        ]))
        return t

    story = [Paragraph(esc(topic) or "Experiment Plan", h_title)]
    n = len(papers or [])
    story.append(Paragraph(
        f"Generated by Sift's Method &amp; Experiment Designer &middot; grounded in "
        f"{n} source{'s' if n != 1 else ''} &middot; {date.today().isoformat()}", small))
    story.append(Spacer(1, 6))
    story.append(HRFlowable(width="100%", color=accent, thickness=1))

    if plan.get("domain"):
        story.append(Paragraph(f"<b>Domain:</b> {esc(plan['domain'])}", body))
    if plan.get("note"):
        story.append(Paragraph(esc(plan["note"]), italic))

    ext_by_idx = {e.get("idx"): e for e in (extractions or [])}
    title_by_idx = {p.get("idx"): p.get("title") for p in (papers or [])}
    crit_by_index = {c.get("index"): c for c in (critique or {}).get("critiques", [])}

    for i, h in enumerate(plan.get("hypotheses") or []):
        story.append(Paragraph(f"H{i + 1}. {esc(h.get('hypothesis'))}", h_hyp))
        if h.get("rationale"):
            story.append(Paragraph(esc(h["rationale"]), italic))

        c = crit_by_index.get(i)
        if c:
            s = c.get("scores") or {}
            story.append(table(
                ["Novelty", "Grounding", "Testability", "Consistency", "Overall"],
                [[s.get("novelty", "—"), s.get("grounding", "—"), s.get("testability", "—"),
                  s.get("consistency", "—"), c.get("overall", "—")]],
            ))
            story.append(Spacer(1, 4))
            if c.get("issues"):
                story.append(Paragraph("Critic: " + esc("; ".join(c["issues"])), small))

        if h.get("setup"):
            story.append(Paragraph(f"<b>Setup:</b> {esc(h['setup'])}", body))
        v = h.get("variables") or {}
        if v:
            story.append(Paragraph(
                f"<b>Variables:</b> independent: {esc(v.get('independent', ''))} &middot; "
                f"dependent: {esc(v.get('dependent', ''))} &middot; "
                f"controlled: {esc(v.get('controlled', ''))}", body))

        grounded, proposed = _evidence_items(h)
        if grounded:
            story.append(Paragraph("Evidence trail", bold_lbl))
            story.append(table(
                ["Approach / baseline", "Source paper", "What it found"],
                [[esc(it["name"]),
                  esc(title_by_idx.get(it["from_idx"])) or f"[{it['from_idx']}]",
                  esc((ext_by_idx.get(it["from_idx"]) or {}).get("finding")) or "—"]
                 for it in grounded],
                col_widths=[1.6 * inch, 2.3 * inch, 2.6 * inch],
            ))
            story.append(Spacer(1, 4))
        if proposed:
            story.append(Paragraph(
                "Not traced to a source (the agent's own proposal): " + esc(", ".join(proposed)), small))

        metrics = h.get("metrics") or []
        if metrics:
            story.append(table(
                ["Metric", "Unit / scale", "Target"],
                [[esc(m.get("name")), esc(m.get("unit")) or "—", esc(m.get("target")) or "—"]
                 for m in metrics],
                col_widths=[2.4 * inch, 1.9 * inch, 2.2 * inch],
            ))
            story.append(Spacer(1, 4))

        if h.get("failure_modes"):
            story.append(Paragraph(f"<b>What could invalidate it:</b> {esc('; '.join(h['failure_modes']))}", body))
        if h.get("validation"):
            story.append(Paragraph(f"<b>Validation:</b> {esc(h['validation'])}", body))
        if h.get("risks"):
            story.append(Paragraph(f"<b>Risks / ethics:</b> {esc(h['risks'])}", body))

    doc.build(story)
    return buf.getvalue()
