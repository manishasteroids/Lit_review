from __future__ import annotations

import io
import re
from html.parser import HTMLParser

import httpx

import ipaddress
import socket
from urllib.parse import urlparse

from core.config import settings

_CACHE: dict[str, str | None] = {}
_PDF_CACHE: dict[str, bytes | None] = {}
MAX_CHARS = 40_000          # ~10k tokens; keeps text prompts affordable
MAX_PDF_BYTES = 25_000_000  # stay under Anthropic's PDF request limit
TIMEOUT = 25.0
HEADERS = {"User-Agent": "Sift-LitReview/1.0 (research assistant)"}


def _is_safe_url(url: str | None) -> bool:
    """SSRF guard: allow only http(s) URLs whose host resolves to a public IP."""
    if not url:
        return False
    try:
        p = urlparse(url)
    except Exception:
        return False
    if p.scheme not in ("http", "https") or not p.hostname:
        return False
    try:
        infos = socket.getaddrinfo(p.hostname, None)
    except Exception:
        return False
    for info in infos:
        try:
            addr = ipaddress.ip_address(info[4][0])
        except ValueError:
            return False
        if (addr.is_private or addr.is_loopback or addr.is_link_local
                or addr.is_reserved or addr.is_multicast or addr.is_unspecified):
            return False
    return True

def _arxiv_pdf_url(url: str) -> str:
    """Turn an arXiv abstract link into its PDF link (full text lives there)."""
    m = re.search(r"arxiv\.org/abs/([\w.\-/]+)", url)
    if m:
        return f"https://arxiv.org/pdf/{m.group(1)}.pdf"
    return url


class _TextExtractor(HTMLParser):
    """Minimal HTML -> visible-text stripper (no external deps)."""

    def __init__(self) -> None:
        super().__init__()
        self._skip = 0
        self.parts: list[str] = []

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style", "noscript", "svg"):
            self._skip += 1

    def handle_endtag(self, tag):
        if tag in ("script", "style", "noscript", "svg") and self._skip:
            self._skip -= 1

    def handle_data(self, data):
        if self._skip == 0:
            t = data.strip()
            if t:
                self.parts.append(t)


def _html_to_text(html: str) -> str:
    p = _TextExtractor()
    try:
        p.feed(html)
    except Exception:
        pass
    return re.sub(r"\n{3,}", "\n\n", "\n".join(p.parts)).strip()


def _pdf_to_text(data: bytes) -> str:
    try:
        from pypdf import PdfReader
    except Exception:
        return ""
    try:
        reader = PdfReader(io.BytesIO(data))
        return "\n".join((page.extract_text() or "") for page in reader.pages).strip()
    except Exception:
        return ""


# Public wrappers for locally-uploaded files (Sources > "Upload a file") —
# same extraction, just named for an external caller instead of the
# URL-fetch pipeline above.
def pdf_bytes_to_text(data: bytes) -> str:
    return _pdf_to_text(data)


def docx_bytes_to_text(data: bytes) -> str:
    try:
        import docx
    except Exception:
        return ""
    try:
        doc = docx.Document(io.BytesIO(data))
        parts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text and cell.text.strip():
                        parts.append(cell.text.strip())
        return "\n".join(parts).strip()
    except Exception:
        return ""


def pptx_bytes_to_text(data: bytes) -> str:
    """Slide text + speaker notes, in slide order — used when a user uploads
    a .pptx as a source (Sources > "Upload a file" / the standalone Studio
    document-upload entry). python-pptx is already a dependency (used for
    the review's own slide-deck export)."""
    try:
        from pptx import Presentation
    except Exception:
        return ""
    try:
        prs = Presentation(io.BytesIO(data))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    slide_parts.append(shape.text_frame.text.strip())
                elif shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                slide_parts.append(cell.text.strip())
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    slide_parts.append(f"[notes] {notes}")
            if slide_parts:
                parts.append(f"[Slide {i}] " + " — ".join(slide_parts))
        return "\n".join(parts).strip()
    except Exception:
        return ""


def _read_local_upload(paper: dict) -> tuple[str, bytes] | None:
    """Read a locally-uploaded source's raw bytes off disk, given the
    `local_file` path set by api/routes.py's upload_paper(). Returns
    (extension, bytes) or None if there's no local file / it's missing."""
    import os

    local_path = (paper.get("local_file") or "").strip()
    if not local_path:
        return None
    full_path = os.path.join(settings.uploads_dir, local_path)
    if not os.path.isfile(full_path):
        return None
    ext = os.path.splitext(full_path)[1].lower()
    try:
        with open(full_path, "rb") as f:
            return ext, f.read()
    except Exception:
        return None


def local_paper_full_text(paper: dict, max_chars: int = MAX_CHARS) -> str | None:
    """Full text for a locally-uploaded paper (Sources > "Upload a file"),
    re-extracted straight from the original file on disk — richer than the
    short excerpt cached on the paper dict at upload time (abstract/_text
    are capped small so the session JSON stays light). This is the
    local-file counterpart to fetch_paper_text(url) above; uploads have no
    URL to fetch, so chat grounding must go through here instead — using
    fetch_paper_text(paper.get("url")) on an uploaded paper silently
    returns nothing, which is why chat answers on uploads used to fall back
    to the truncated abstract only."""
    read = _read_local_upload(paper)
    if not read:
        return None
    ext, data = read
    if ext == ".pdf":
        text = _pdf_to_text(data)
    elif ext == ".pptx":
        text = pptx_bytes_to_text(data)
    elif ext == ".docx":
        text = docx_bytes_to_text(data)
    else:
        return None
    text = (text or "").strip()
    if not text:
        return None
    return text[:max_chars] + "\n\n[...truncated...]" if len(text) > max_chars else text


def local_paper_pdf_bytes(paper: dict) -> bytes | None:
    """Raw PDF bytes for a locally-uploaded paper, straight from disk — the
    local-file counterpart to fetch_paper_pdf(url) above. Lets chat attach
    the real PDF (so the model can read figures/diagrams/tables visually)
    for an uploaded source, the same way it already can for a URL-backed
    one; without this, a diagram question on an uploaded paper silently got
    no PDF attached at all since fetch_paper_pdf(None) always returns None."""
    read = _read_local_upload(paper)
    if not read:
        return None
    ext, data = read
    if ext != ".pdf" or len(data) > MAX_PDF_BYTES:
        return None
    return data


_DOI_RE = re.compile(r"10\.\d{4,9}/[^\s\"'<>&]+", re.I)


def _fetch_pdf_bytes(url: str) -> bytes | None:
    """GET a URL and return the bytes only if it's actually a PDF."""
    if not _is_safe_url(url):
        return None
    try:
        with httpx.Client(follow_redirects=True, timeout=TIMEOUT, headers=HEADERS) as client:
            resp = client.get(url)
            resp.raise_for_status()
            ctype = resp.headers.get("content-type", "").lower()
            if "pdf" in ctype or url.lower().split("?")[0].endswith(".pdf"):
                return resp.content
    except Exception:
        return None
    return None


def _extract_doi(url: str) -> str | None:
    m = _DOI_RE.search(url or "")
    return m.group(0).rstrip(").,;") if m else None


def _unpaywall_pdf_url(doi: str) -> str | None:
    """Ask Unpaywall for an open-access PDF for a DOI (any repository / PMC / publisher OA)."""
    email = getattr(settings, "unpaywall_email", "") or "research@example.com"
    try:
        with httpx.Client(follow_redirects=True, timeout=TIMEOUT, headers=HEADERS) as client:
            r = client.get(f"https://api.unpaywall.org/v2/{doi}", params={"email": email})
            r.raise_for_status()
            j = r.json()
    except Exception:
        return None
    best = j.get("best_oa_location") or {}
    if best.get("url_for_pdf"):
        return best["url_for_pdf"]
    for loc in (j.get("oa_locations") or []):
        if loc.get("url_for_pdf"):
            return loc["url_for_pdf"]
    return None


def fetch_paper_pdf(url: str | None) -> bytes | None:
    """Return the raw PDF bytes for a paper (open-access), or None.

    Tries the link directly (handles arXiv abs→pdf and any .pdf); if that isn't
    a PDF, resolves an open-access copy via Unpaywall using the DOI. Lets the
    chat read the WHOLE paper (figures + tables) for far more papers, falling
    back to text/abstract only when no OA PDF exists anywhere."""
    if not url or not url.startswith(("http://", "https://")):
        return None
    if url in _PDF_CACHE:
        return _PDF_CACHE[url]

    data = _fetch_pdf_bytes(_arxiv_pdf_url(url))

    if not data:
        doi = _extract_doi(url)
        if doi:
            oa = _unpaywall_pdf_url(doi)
            if oa:
                data = _fetch_pdf_bytes(oa)

    if data and len(data) > MAX_PDF_BYTES:
        data = None
    _PDF_CACHE[url] = data
    return data


def fetch_paper_text(url: str | None) -> str | None:
    """Return extracted full text for a paper URL, or None if unavailable."""
    if not url or not url.startswith(("http://", "https://")):
        return None
    if url in _CACHE:
        return _CACHE[url]

    target = _arxiv_pdf_url(url)
    if not _is_safe_url(target):
        _CACHE[url] = None
        return None
    text: str | None = None
    try:
        with httpx.Client(follow_redirects=True, timeout=TIMEOUT, headers=HEADERS) as client:
            resp = client.get(target)
            resp.raise_for_status()
            ctype = resp.headers.get("content-type", "").lower()
            if "pdf" in ctype or target.lower().endswith(".pdf"):
                text = _pdf_to_text(resp.content)
            else:
                text = _html_to_text(resp.text)
    except Exception:
        text = None

    if text:
        text = text.strip()
        if len(text) > MAX_CHARS:
            text = text[:MAX_CHARS] + "\n\n[...truncated...]"
    result = text or None
    _CACHE[url] = result
    return result
